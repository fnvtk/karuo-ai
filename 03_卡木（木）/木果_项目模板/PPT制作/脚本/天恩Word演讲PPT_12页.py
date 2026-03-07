#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
天恩 Word（笛卡尔背单词）项目 · 小学生演讲 PPT 生成
汇报人：二年级4班 施吴佶 | 12页 | 小学生口吻 | 含流程图、插图、讲解
风格：仿天恩乖乖，暖色、圆角、图标、童趣；每页有图或流程图。
"""
from pathlib import Path
import sys

# 输出目录：优先天恩项目下，便于用户使用
TIANEN_ROOT = Path("/Users/karuo/Library/Mobile Documents/com~apple~CloudDocs/Documents/婼瑄/天恩")
OUT_DIR = TIANEN_ROOT / "天恩Word演讲PPT"
KARUO_IMG = Path("/Users/karuo/Documents/卡若Ai的文件夹/图片")

# 确保可导入 pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("请先安装: pip install python-pptx")
    sys.exit(1)


# ---------- 配色（小学绘本/童趣风 + 苹果毛玻璃感）----------
BG_CREAM = RGBColor(255, 250, 230)      # 暖黄
BG_PINK = RGBColor(255, 245, 248)       # 浅粉
BG_GLASS = RGBColor(248, 248, 252)      # 毛玻璃浅白
TITLE_BROWN = RGBColor(139, 90, 43)     # 标题棕
TEXT_DARK = RGBColor(60, 45, 30)       # 正文深棕
ACCENT = RGBColor(218, 165, 32)         # 金黄强调
BORDER = RGBColor(180, 140, 100)       # 边框暖棕
GLASS_BORDER = RGBColor(220, 220, 230)  # 毛玻璃边框
LIGHT_PINK = RGBColor(255, 228, 225)   # 封面渐变用


def set_slide_background(slide, color=BG_CREAM):
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    except Exception:
        pass


def add_title(slide, text, top=Inches(0.4), left=Inches(0.6), width=Inches(12), font_size=28, center=False):
    """添加标题；center=True 时文字居中"""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = TITLE_BROWN
    if center:
        p.alignment = PP_ALIGN.CENTER
    return tb


def add_body(slide, lines, top, left=Inches(0.6), width=Inches(7), font_size=20, center=False):
    tb = slide.shapes.add_textbox(left, top, width, Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)
        if center:
            p.alignment = PP_ALIGN.CENTER
    return tb


def add_flowchart(slide, top=Inches(2), left=Inches(0.8), glass_style=True):
    """绘制流程图（毛玻璃风格）：选后缀 → 加前缀 → 生成单词 → 听句子看图 → 标记已背"""
    box_w, box_h = Inches(1.85), Inches(0.65)
    steps = ["1. 选后缀", "2. 加前缀", "3. 生成单词", "4. 听句看图", "5. 标记已背"]
    fill_rgb = BG_GLASS if glass_style else RGBColor(255, 255, 255)
    line_rgb = GLASS_BORDER if glass_style else BORDER
    for i, label in enumerate(steps):
        x = left + Inches(i * 2.15)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            arrow_x = x + box_w + Inches(0.04)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, top + Inches(0.18), Inches(0.3), Inches(0.28)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.color.rgb = ACCENT


SLIDE_W = Inches(13.333)

def add_picture_glass(slide, path, left, top, width, border_rgb=GLASS_BORDER, center=False):
    """插入图片并加毛玻璃感边框。center=True 时水平居中。"""
    path = Path(path)
    if not path.exists():
        return None
    try:
        if center:
            left = (SLIDE_W - width) / 2
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
        pic.line.color.rgb = border_rgb
        pic.line.width = Pt(2)
        return pic
    except Exception:
        return None


def add_cartesian_diagram(slide, top=Inches(2.2), left=Inches(0.8), glass_style=True):
    """笛卡尔坐标简易示意图（毛玻璃风格）：X轴后缀、Y轴前缀、交叉得单词"""
    fill_rgb = BG_GLASS if glass_style else RGBColor(255, 255, 255)
    line_rgb = GLASS_BORDER if glass_style else BORDER
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(8), Inches(2.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_rgb
    box.line.color.rgb = line_rgb
    box.line.width = Pt(1.5)
    # X轴标签
    tb_x = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), Inches(2), Inches(0.4))
    tb_x.text_frame.paragraphs[0].text = "X轴：后缀 (DAY, BOOK...)"
    tb_x.text_frame.paragraphs[0].font.size = Pt(16)
    tb_x.text_frame.paragraphs[0].font.bold = True
    tb_x.text_frame.paragraphs[0].font.color.rgb = TITLE_BROWN
    # Y轴标签
    tb_y = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1), Inches(2), Inches(0.4))
    tb_y.text_frame.paragraphs[0].text = "Y轴：前缀 (SUN, NOTE...)"
    tb_y.text_frame.paragraphs[0].font.size = Pt(16)
    tb_y.text_frame.paragraphs[0].font.bold = True
    tb_y.text_frame.paragraphs[0].font.color.rgb = TITLE_BROWN
    # 示例单词
    tb_w = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.8), Inches(7), Inches(0.8))
    tb_w.text_frame.word_wrap = True
    p = tb_w.text_frame.paragraphs[0]
    p.text = "组合起来 → SUNDAY, MONDAY, NOTEBOOK, BIRTHDAY ... 一次记住好多词！"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_DARK


def add_icon_text(slide, icon, text, left, top, font_size=22):
    tb = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{icon}  {text}"
    p.font.size = Pt(font_size)
    p.font.color.rgb = TEXT_DARK
    return tb


def create_placeholder_photo(path, name="施吴佶", size=400):
    """生成一张占位图：渐变底 + 圆形 + 姓名，用于演讲人照片位"""
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return None
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (size, size), (255, 240, 230))
    draw = ImageDraw.Draw(img)
    # 渐变
    for y in range(size):
        t = y / size
        r = int(255)
        g = int(240 + (220 - 240) * t)
        b = int(230 + (200 - 230) * t)
        draw.line([(0, y), (size, y)], fill=(r, g, b))
    # 圆
    margin = 60
    draw.ellipse([margin, margin, size - margin, size - margin], outline=(180, 140, 100), width=4)
    # 姓名
    try:
        from PIL import ImageFont
        font = ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", 36)
    except Exception:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), name, font=font)
    tw = bbox[2] - bbox[0]
    draw.text(((size - tw) // 2, size // 2 - 25), name, font=font, fill=(139, 90, 43))
    img.save(path, "PNG")
    return path


def build_presentation(photo_path=None, extra_images=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 准备演讲人照片与插图资源
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    assets_dir = OUT_DIR / "assets"
    speaker1 = assets_dir / "speaker1.png"  # 小女孩剪刀手
    speaker2 = assets_dir / "speaker2.png"  # 户外照
    coordinate_img = assets_dir / "coordinate.png"
    flow_words_img = assets_dir / "flow_words.png"
    scale_img = assets_dir / "scale.png"
    sunflower_book_img = assets_dir / "sunflower_book.png"
    day_table_img = assets_dir / "day_table.png"

    speaker_img = photo_path
    if not speaker_img or not Path(speaker_img).exists():
        if speaker1.exists():
            speaker_img = str(speaker1)
        else:
            placeholder_photo = OUT_DIR / "placeholder_speaker.png"
            created = create_placeholder_photo(placeholder_photo)
            speaker_img = str(placeholder_photo) if (created and placeholder_photo.exists()) else None
    else:
        speaker_img = str(speaker_img)
    speaker2_path = str(speaker2) if speaker2.exists() else speaker_img

    def add_slide():
        return prs.slides.add_slide(blank)

    # ---------- 第1页：封面（标题居中、图片居中放大、汇报人一行）----------
    s1 = add_slide()
    set_slide_background(s1, BG_CREAM)
    add_title(s1, "天恩 Word · 智能单词记忆小助手", top=Inches(0.7), left=Inches(0.5), width=Inches(12.3), font_size=40, center=True)
    add_body(s1, ["向全班、全校、全国的小朋友介绍我们的项目"], top=Inches(1.4), left=Inches(1.5), width=Inches(10.3), font_size=22, center=True)
    if speaker_img and Path(speaker_img).exists():
        add_picture_glass(s1, speaker_img, 0, Inches(2.0), Inches(3.2), BORDER, center=True)
    add_title(s1, "汇报人：二年四班 4号 施吴佶", top=Inches(5.6), left=Inches(0.5), width=Inches(12.3), font_size=24, center=True)

    # ---------- 第2页：大家好，我是施吴佶（双图可用 speaker2）----------
    s2 = add_slide()
    set_slide_background(s2, BG_CREAM)
    add_title(s2, "大家好，我是施吴佶！", font_size=32, center=True)
    add_icon_text(s2, "👋", "我是二年级4班的一名小学生。", Inches(0.6), Inches(1.4))
    add_icon_text(s2, "📚", "我和小伙伴们一起做了一个「记单词」的小项目，", Inches(0.6), Inches(2.0))
    add_icon_text(s2, "✨", "今天想跟大家说一说它是怎么用的！", Inches(0.6), Inches(2.6))
    if speaker2_path and Path(speaker2_path).exists():
        add_picture_glass(s2, speaker2_path, Inches(7.8), Inches(1.6), Inches(2.6), BORDER)

    # ---------- 第3页：今天讲什么 ----------
    s3 = add_slide()
    set_slide_background(s3, BG_CREAM)
    add_title(s3, "今天我要讲什么？", font_size=30, center=True)
    add_body(s3, [
        "我们要讲的是一个「智能单词记忆小助手」！",
        "它有一个很厉害的方法，叫「笛卡尔坐标记忆法」。",
        "用这个方法，可以一次记住好多相关的单词，又好玩又不容易忘。",
        "后面我会告诉大家怎么用、我们班有什么收获。"
    ], top=Inches(1.5), font_size=20)

    # ---------- 第4页：为什么做这个（可配情绪/改进示意图）----------
    s4 = add_slide()
    set_slide_background(s4, BG_PINK)
    add_title(s4, "我们为什么做这个？", font_size=30, center=True)
    add_icon_text(s4, "😣", "记单词太难了！", Inches(0.6), Inches(1.3), 24)
    add_body(s4, [
        "以前记单词总是死记硬背，很容易忘。",
        "而且一个一个记，又慢又没意思。",
        "我们就想：能不能把有关的词放在一起记？",
        "所以做了这个小助手，让记单词变得更有趣、更快。"
    ], top=Inches(1.9), left=Inches(0.6), width=Inches(6.2), font_size=20)
    if scale_img.exists():
        add_picture_glass(s4, scale_img, Inches(7.5), Inches(1.8), Inches(4.8), BORDER)

    # ---------- 第5页：我们的办法 — 笛卡尔坐标 + 坐标系插图 ----------
    s5 = add_slide()
    set_slide_background(s5, BG_CREAM)
    add_title(s5, "我们的办法：笛卡尔坐标记忆法", font_size=28, center=True)
    add_body(s5, [
        "把单词拆成「前缀」和「后缀」：",
        "X 轴放后缀（比如 DAY、BOOK），Y 轴放前缀（比如 SUN、NOTE）。",
        "它们一组合，就变成 SUNDAY、NOTEBOOK 这些词，一次能记一串！"
    ], top=Inches(1.0), left=Inches(0.6), width=Inches(5.8), font_size=19)
    if coordinate_img.exists():
        add_picture_glass(s5, coordinate_img, Inches(6.8), Inches(2.0), Inches(5.5), GLASS_BORDER)
    else:
        add_cartesian_diagram(s5, top=Inches(2.9), left=Inches(0.6))

    # ---------- 第6页：什么是笛卡尔坐标 + DAY 单词表示例图 ----------
    s6 = add_slide()
    set_slide_background(s6, BG_CREAM)
    add_title(s6, "什么是笛卡尔坐标？", font_size=30, center=True)
    add_body(s6, [
        "就像画一个「十字」，横着的是后缀，竖着的是前缀。",
        "它们交叉的地方，就是一个新单词！",
        "这样记，单词和单词之间就有联系，不会乱。"
    ], top=Inches(1.2), left=Inches(0.6), width=Inches(6), font_size=20)
    if day_table_img.exists():
        add_picture_glass(s6, day_table_img, Inches(7.0), Inches(1.8), Inches(5.5), GLASS_BORDER)
    else:
        add_cartesian_diagram(s6, top=Inches(3.2), left=Inches(1.5))

    # ---------- 第7页：怎么用 — 流程图（毛玻璃）+ 步骤示意图 ----------
    s7 = add_slide()
    set_slide_background(s7, BG_PINK)
    add_title(s7, "怎么用？五步就会！", font_size=30, center=True)
    add_flowchart(s7, top=Inches(1.85), left=Inches(0.5), glass_style=True)
    if flow_words_img.exists():
        add_picture_glass(s7, flow_words_img, Inches(7.2), Inches(2.0), Inches(5.3), GLASS_BORDER)
    add_body(s7, [
        "选一个后缀 → 加上几个前缀 → 点「生成」→ 听句子、看配图 → 会了就标记「已背」！"
    ], top=Inches(3.5), left=Inches(0.6), width=Inches(12), font_size=18)

    # ---------- 第8页：记忆句和配图 + 联想示意图 ----------
    s8 = add_slide()
    set_slide_background(s8, BG_CREAM)
    add_title(s8, "记忆句和配图", font_size=30, center=True)
    add_icon_text(s8, "📝", "一句子里藏着这一组的所有单词，读一句就复习一遍。", Inches(0.6), Inches(1.3))
    add_icon_text(s8, "🖼️", "还有 AI 画的配图，看图就能想起这句话和单词。", Inches(0.6), Inches(1.9))
    add_body(s8, [
        "这样记单词又轻松又牢，我们班同学都很喜欢！"
    ], top=Inches(2.5), left=Inches(0.6), width=Inches(5.8), font_size=20)
    if sunflower_book_img.exists():
        add_picture_glass(s8, sunflower_book_img, Inches(6.6), Inches(1.5), Inches(6), GLASS_BORDER)

    # ---------- 第9页：界面长什么样 + 单词表示例 ----------
    s9 = add_slide()
    set_slide_background(s9, BG_CREAM)
    add_title(s9, "我们做的界面长什么样？", font_size=28, center=True)
    add_body(s9, [
        "打开网页就能用，不用装软件。",
        "上面是「坐标系」：左边是前缀，上边是后缀，中间是单词和意思。",
        "下面有一句「记忆句」和一张图，点一下还能朗读。",
        "会了的词点一下就能标成「已背」，特别方便！"
    ], top=Inches(1.2), left=Inches(0.6), width=Inches(6), font_size=19)
    if day_table_img.exists():
        add_picture_glass(s9, day_table_img, Inches(6.8), Inches(1.8), Inches(5.5), GLASS_BORDER)

    # ---------- 第10页：我们班的收获 ----------
    s10 = add_slide()
    set_slide_background(s10, BG_PINK)
    add_title(s10, "我们班的收获", font_size=30, center=True)
    add_icon_text(s10, "🌟", "记单词变快了，而且不容易忘。", Inches(0.6), Inches(1.4))
    add_icon_text(s10, "😊", "大家觉得好玩，都愿意多背几组。", Inches(0.6), Inches(2.0))
    add_icon_text(s10, "🤝", "我们一起想点子、一起用，更像一个小组了。", Inches(0.6), Inches(2.6))
    add_body(s10, ["希望更多班级、更多小朋友也能用上，一起轻松记单词！"], top=Inches(3.2), font_size=20)

    # ---------- 第11页：谢谢大家 ----------
    s11 = add_slide()
    set_slide_background(s11, BG_CREAM)
    add_title(s11, "谢谢大家！", top=Inches(2.2), font_size=44, center=True)
    add_body(s11, ["欢迎大家一起用「天恩 Word」记单词～"], top=Inches(3.0), font_size=24)
    add_body(s11, ["汇报人：二年四班 4号 施吴佶"], top=Inches(4.0), font_size=20, center=True)

    # ---------- 第12页：感谢 ----------
    s12 = add_slide()
    set_slide_background(s12, LIGHT_PINK)
    add_title(s12, "感谢聆听", top=Inches(2.8), font_size=38, center=True)
    add_body(s12, ["感谢全班、全校、全国的老师和小伙伴们！"], top=Inches(3.6), font_size=22)

    out_ppt = OUT_DIR / "天恩Word项目介绍_二年级4班施吴佶_12页.pptx"
    prs.save(out_ppt)
    print("已生成:", out_ppt)
    return out_ppt


def main():
    import argparse
    ap = argparse.ArgumentParser(description="天恩 Word 小学生演讲 PPT（12页）")
    ap.add_argument("--photo", default="", help="施吴佶照片路径，无则用占位图")
    ap.add_argument("--output-dir", default="", help="输出目录，默认天恩项目下")
    args = ap.parse_args()
    if args.output_dir:
        global OUT_DIR
        OUT_DIR = Path(args.output_dir)
    build_presentation(photo_path=args.photo or None)


if __name__ == "__main__":
    main()
