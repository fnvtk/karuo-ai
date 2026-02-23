#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
天恩和乖乖 - 二年级绘本风格 PPT（优化版）
参考：我的植物故事、我的水果拼盘之旅 - 黄色底、图片有边框、每图配文字
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

DIR_BASE = Path("/Users/karuo/Library/Mobile Documents/com~apple~CloudDocs/Documents/婼瑄/天恩/ppt 202060223")
OUT_PATH = DIR_BASE / "我和乖乖的故事.pptx"

# 风格配色（参考黄色底绘本）
BG_YELLOW = RGBColor(255, 250, 230)      # 暖黄背景 #FFFAE6
TITLE_BROWN = RGBColor(139, 90, 43)      # 标题棕色
TEXT_DARK = RGBColor(60, 45, 30)         # 正文深色
BORDER_BROWN = RGBColor(180, 140, 100)   # 图片边框暖棕
ACCENT = RGBColor(218, 165, 32)          # 强调色金黄


def set_slide_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_picture_with_border(slide, img_path, left, top, width, caption=None):
    """添加带边框和说明的图片"""
    pic = slide.shapes.add_picture(str(img_path), left, top, width=width)
    pic.line.color.rgb = BORDER_BROWN
    pic.line.width = Pt(3)
    if caption:
        cap_top = top + pic.height + Pt(6)  # 图片下方留小间距
        tb = slide.shapes.add_textbox(left, cap_top, width, Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📷 {caption}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.alignment = 1
    return pic


def ensure_images():
    imgs = sorted(DIR_BASE.glob("*.jpg"))
    return [str(p) for p in imgs]


def add_title_slide(prs):
    """第1页：封面 - 黄色底"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_background(slide, BG_YELLOW)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.3), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "我和乖乖的故事"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_BROWN
    p.alignment = 1
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "汇报人：天恩"
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEXT_DARK
    p2.alignment = 1
    p3 = tf2.add_paragraph()
    p3.text = "班级：二年级4班"
    p3.font.size = Pt(24)
    p3.font.color.rgb = TEXT_DARK
    p3.alignment = 1


def add_content_slide(prs, title, body_lines, img_info=None):
    """内容页：标题+正文+带边框图片+图片说明
    img_info: [(path, caption), ...] 或 [(path, caption)]
    """
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_background(slide, BG_YELLOW)
    left = Inches(0.5)
    top = Inches(0.4)
    tb = slide.shapes.add_textbox(left, top, Inches(12.3), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_BROWN
    body_top = Inches(1.2)
    if img_info and len(img_info) > 0:
        txt_width = Inches(6)
    else:
        txt_width = Inches(12.3)
    tb2 = slide.shapes.add_textbox(left, body_top, txt_width, Inches(5.2))
    tf = tb2.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)
    if img_info and len(img_info) >= 1:
        first = img_info[0]
        img_path, caption = (first[0], first[1]) if isinstance(first, (list, tuple)) else (first, "")
        if os.path.exists(img_path):
            if len(img_info) == 2:
                add_picture_with_border(slide, img_path, Inches(6.6), Inches(1.2), Inches(2.8), caption)
                sec = img_info[1]
                img_path2, caption2 = (sec[0], sec[1]) if isinstance(sec, (list, tuple)) else (sec, "")
                if os.path.exists(img_path2):
                    add_picture_with_border(slide, img_path2, Inches(9.6), Inches(1.2), Inches(2.8), caption2)
            else:
                add_picture_with_border(slide, img_path, Inches(6.6), Inches(1.2), Inches(5.8), caption)


def add_story_slide(prs, title, story_lines, img_info=None):
    """故事页：左文右图，图带边框和说明"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_background(slide, BG_YELLOW)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(26)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = TITLE_BROWN
    txt_width = Inches(6) if img_info else Inches(12.3)
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), txt_width, Inches(5.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    for i, line in enumerate(story_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)
    if img_info:
        img_path, caption = img_info
        if os.path.exists(img_path):
            add_picture_with_border(slide, img_path, Inches(6.6), Inches(1.2), Inches(5.8), caption)


def main():
    imgs = ensure_images()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs)

    # 第2页：这是我和乖乖！
    add_content_slide(
        prs,
        "这是我和乖乖！",
        [
            "乖乖是我家的小狗，它是我的好朋友。",
            "每天放学回家，乖乖都会在门口等我，摇着尾巴欢迎我。",
            "我们一起玩耍、一起散步，乖乖让我的生活变得更快乐！",
        ],
        [(imgs[0], "这张图：我和乖乖的合影")] if imgs else None,
    )

    # 第3页：乖乖是怎么来到我家的？
    add_story_slide(
        prs,
        "乖乖是怎么来到我家的？",
        [
            "乖乖是爸爸妈妈送给我的惊喜。",
            "第一次见到它的时候，它小小的、毛茸茸的，一双黑溜溜的眼睛望着我。",
            "从那天起，我们就成了形影不离的好伙伴！",
        ],
        (imgs[1], "乖乖刚来我家的样子") if len(imgs) > 1 else None,  # 单图用元组
    )

    # 第4页：我和乖乖每天做什么？
    add_story_slide(
        prs,
        "我和乖乖每天做什么？",
        [
            "早上，乖乖会叫醒我起床。",
            "放学后，我们一起在小区里跑步、追球。",
            "乖乖最喜欢我摸它的脑袋，每次摸它，它都会舒服地眯起眼睛。",
            "我们还会一起看绘本，乖乖趴在我身边，好像在听故事呢！",
        ],
        (imgs[2], "我和乖乖一起玩") if len(imgs) > 2 else None,  # 单图用元组
    )

    # 第5页：乖乖最可爱的地方
    add_content_slide(
        prs,
        "乖乖最可爱的地方",
        [
            "乖乖的耳朵软软的，跑起来一甩一甩的。",
            "它特别聪明，听得懂「坐下」「握手」这些口令。",
            "乖乖还很懂事，我写作业的时候它会安静地陪着我。",
        ],
        [
            (imgs[3], "乖乖的可爱瞬间"),
            (imgs[4], "乖乖在等我"),
        ] if len(imgs) > 4 else ([(imgs[3], "乖乖的可爱瞬间")] if len(imgs) > 3 else None),
    )

    # 第6页：更多和乖乖的快乐时光
    add_content_slide(
        prs,
        "更多和乖乖的快乐时光",
        [
            "我们一起出门散步，乖乖总是跑在我前面。",
            "下雨天，乖乖会趴在窗边陪我一起看雨。",
            "乖乖是我最好的朋友，有它陪伴的每一天都很开心！",
        ],
        [(imgs[5], "散步时光")] if len(imgs) > 5 else None,
    )

    # 第7页：我为乖乖写了一首小诗
    add_content_slide(
        prs,
        "我为乖乖写了一首小诗",
        [
            "《我的小伙伴乖乖》",
            "软软的耳朵，黑黑的眼睛，",
            "乖乖是我最好的朋友。",
            "一起奔跑，一起玩耍，",
            "有乖乖的日子，每天都很快乐！",
        ],
        [(imgs[6], "乖乖和我的日常")] if len(imgs) > 6 else None,
    )

    # 第8页：更多图片展示
    if len(imgs) > 7:
        add_content_slide(
            prs,
            "我们的珍贵回忆",
            [
                "这些照片记录了我和乖乖在一起的美好时光。",
                "每一张都是我们友谊的见证！",
            ],
            [
                (imgs[7], "回忆一"),
                (imgs[8], "回忆二"),
            ] if len(imgs) > 8 else [(imgs[7], "珍贵回忆")],
        )

    # 结尾页
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_background(slide, BG_YELLOW)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.3), Inches(2.5))
    tf = tb.text_frame
    tf.paragraphs[0].text = "我的故事讲完啦，希望大家喜欢！"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TITLE_BROWN
    tf.paragraphs[0].alignment = 1
    p2 = tf.add_paragraph()
    p2.text = "你们家里也有小动物好朋友吗？"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_DARK
    p2.alignment = 1
    p3 = tf.add_paragraph()
    p3.text = "谢谢大家！"
    p3.font.size = Pt(26)
    p3.font.bold = True
    p3.font.color.rgb = TITLE_BROWN
    p3.alignment = 1

    prs.save(str(OUT_PATH))
    print("✅ PPT 已优化生成:", OUT_PATH)


if __name__ == "__main__":
    main()
