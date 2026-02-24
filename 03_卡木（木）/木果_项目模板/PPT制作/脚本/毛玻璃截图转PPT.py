#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
毛玻璃风格 HTML → 截图 → 组装 PPT
依赖：playwright（pip install playwright && playwright install chromium）
"""
import subprocess
import sys
from pathlib import Path

# 尝试 playwright，若无则退回「仅组装」模式
try:
    from playwright.sync_api import sync_playwright
    HAS_PLAYWRIGHT = True
except ImportError:
    HAS_PLAYWRIGHT = False

from pptx import Presentation
from pptx.util import Inches

BASE = Path(__file__).resolve().parent
OUT_ROOT = Path("/Users/karuo/Documents/卡若Ai的文件夹/报告")
TIANEN_DIR = Path("/Users/karuo/Library/Mobile Documents/com~apple~CloudDocs/Documents/婼瑄/天恩/ppt 202060223")


def screenshot_slides(html_path, out_slides_dir, max_slides=10):
    """用 playwright 截取每页"""
    if not HAS_PLAYWRIGHT:
        print("⚠️ playwright 未安装")
        return []
    out_slides_dir.mkdir(parents=True, exist_ok=True)
    imgs = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720})
        page.goto(f"file://{html_path}")
        for i in range(1, max_slides + 1):
            sel = f"#slide-{i}"
            el = page.locator(sel)
            if el.count() > 0:
                path = out_slides_dir / f"slide_{i:02d}.png"
                el.screenshot(path=path)
                imgs.append(str(path))
        browser.close()
    return imgs


def build_ppt(imgs, out_ppt):
    """将图片组装成 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    for p in imgs:
        if Path(p).exists():
            s = prs.slides.add_slide(layout)
            s.shapes.add_picture(p, Inches(0), Inches(0), width=prs.slide_width)
    out_ppt.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_ppt)
    return out_ppt


def main():
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--html", default="复盘", choices=["复盘", "卡若人设", "纳瓦尔访谈", "天恩乖乖", "今日日志总结", "家里NAS对话描述", "公司财务分析", "公司财务月报_2026-01"])
    args = ap.parse_args()
    if args.html == "卡若人设":
        html = BASE / "卡若人设PPT_毛玻璃.html"
        out_slides = OUT_ROOT / "卡若人设_毛玻璃_slides"
        out_ppt = OUT_ROOT / "卡若人设介绍_毛玻璃.pptx"
        max_slides = 5
    elif args.html == "家里NAS对话描述":
        html = BASE / "家里NAS对话描述PPT_毛玻璃.html"
        out_slides = OUT_ROOT / "家里NAS对话描述_毛玻璃_slides"
        out_ppt = OUT_ROOT / "家里NAS_对话描述_毛玻璃.pptx"
        max_slides = 8
    elif args.html == "纳瓦尔访谈":
        html = BASE / "纳瓦尔访谈PPT_毛玻璃.html"
        # v2：扩展页数（含方法/问答/流程图/行动清单），避免覆盖旧版
        out_slides = OUT_ROOT / "纳瓦尔访谈_毛玻璃_slides_v2"
        out_ppt = OUT_ROOT / "纳瓦尔访谈_读书笔记_毛玻璃_v2.pptx"
        max_slides = 15
    elif args.html == "公司财务分析":
        html = BASE / "公司财务分析PPT_毛玻璃.html"
        out_slides = OUT_ROOT / "公司财务分析_毛玻璃_slides"
        out_ppt = OUT_ROOT / "公司财务_多维分析_CFO大白话_毛玻璃.pptx"
        max_slides = 12
    elif args.html == "公司财务月报_2026-01":
        html = BASE / "公司财务月报_2026年1月PPT_毛玻璃.html"
        out_slides = OUT_ROOT / "公司财务月报_2026-01_毛玻璃_slides"
        out_ppt = OUT_ROOT / "公司财务月报_2026年1月_CFO大白话_毛玻璃.pptx"
        max_slides = 10
    elif args.html == "天恩乖乖":
        html = BASE / "天恩乖乖PPT_毛玻璃.html"
        out_slides = TIANEN_DIR / "乖乖_毛玻璃_slides"
        out_ppt = TIANEN_DIR / "我和乖乖的故事_高级版.pptx"
        max_slides = 10
    elif args.html == "今日日志总结":
        html = BASE / "今日日志阅读总结PPT_毛玻璃.html"
        out_slides = OUT_ROOT / "今日日志总结_毛玻璃_slides"
        out_ppt = OUT_ROOT / "今日日志阅读总结_毛玻璃.pptx"
        max_slides = 6
    else:
        html = BASE / "复盘PPT_毛玻璃.html"
        out_slides = OUT_ROOT / "复盘_毛玻璃_slides"
        out_ppt = OUT_ROOT / "复盘_2026-02-23_毛玻璃.pptx"
        max_slides = 5
    if HAS_PLAYWRIGHT:
        imgs = screenshot_slides(html, out_slides, max_slides)
    else:
        imgs = sorted([str(p) for p in out_slides.glob("slide_*.png")])
    if not imgs:
        print("❌ 无截图可用。HTML 路径:", html)
        sys.exit(1)
    build_ppt(imgs, out_ppt)
    print("✅ PPT 已生成:", out_ppt)
    # 导出后自动打开输出文件夹，方便直接访问
    import subprocess
    open_dir = str(TIANEN_DIR) if args.html == "天恩乖乖" else str(OUT_ROOT)
    subprocess.run(["open", open_dir], check=False)


if __name__ == "__main__":
    main()
